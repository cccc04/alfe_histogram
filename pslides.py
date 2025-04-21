from pptx import Presentation
from pptx.util import Inches
import os
import pandas as pd

# Function to load parameters from the Excel file
def load_parameters(csv_file):
    df = pd.read_csv(csv_file)
    parameter = df['parameter name']
    parameters = parameter.dropna().astype(str).tolist()

    # Sort parameters by length (longest first) to avoid substring matching issues
    parameters.sort(key=len, reverse=True)

    print(f"Loaded parameters: {parameters}")  # Debugging: Print loaded parameters
    return parameters

# Function to load image files and categorize them by parameter, HG/LG, and 25/50
def load_images(image_folder, parameters):
    images = {}
    matched_images = set()  # Track already matched images

    # Sort parameters by length (longest first)
    sorted_parameters = sorted(parameters, key=lambda param: len(param), reverse=True)

    # Move parameters containing 'uniformity' to the front while maintaining length-based sorting
    uniformity_params = [param for param in sorted_parameters if "uniformity" in param.lower()]
    non_uniformity_params = [param for param in sorted_parameters if "uniformity" not in param.lower()]

    # Combine them: uniformity parameters first, then the rest
    sorted_parameters = uniformity_params + non_uniformity_params

    # Iterate through all files in the folder
    for filename in os.listdir(image_folder):
        if filename.endswith('.png'):
            # Sort parameters by length to match the longest first (specific parameters first)
            matching_parameters = [param for param in sorted_parameters if param in filename]

            if not matching_parameters:
                continue  # Skip files that don't match any parameter
            
            # Assume the longest match is the correct one (since it's sorted by length)
            parameter = matching_parameters[0]
            
            # Skip the image if it has already been matched
            if filename in matched_images and 'uniformity' not in filename:
                continue

            # Mark this image as matched
            matched_images.add(filename)

            # Check for HG/LG in the filename
            if 'HG' in filename or 'hg' in filename:
                hg_lg = 'HG'
            elif 'LG' in filename or 'lg' in filename:
                hg_lg = 'LG'
            else:
                hg_lg = None

            # Check for the suffix (25 or 50)
            if '25' in filename:
                suffix = '25'
            elif '50' in filename:
                suffix = '50'
            else:
                suffix = None  # Skip if neither 25 nor 50 is found in the filename

            # Initialize the dictionary for the images if not already present
            if parameter not in images:
                images[parameter] = {
                    'HG': {'25': [], '50': [], 'None': []}, 
                    'LG': {'25': [], '50': [], 'None': []}, 
                    'None': {'25': [], '50': [], 'None': []}
                }

            # Categorize the image
            image_path = os.path.join(image_folder, filename)
            if not os.path.exists(image_path):
                print(f"Image not found: {image_path}")  # Debugging: check if image exists
                continue

            if hg_lg and suffix:
                images[parameter][hg_lg][suffix].append(image_path)
            elif suffix:
                images[parameter]['None'][suffix].append(image_path)
            else:
                images[parameter]['None']['None'].append(image_path)

    print(f"Categorized images: {images}")  # Debugging: print the categorized images
    return images

def create_presentation(images, output_pptx):
    prs = Presentation()
    prs.slide_width = Inches(13.33)

    # Adding images to slides
    def add_images_to_slide(slide, images, col):
        left = Inches(0.05)  # Adjust left padding to remove extra space
        top = Inches(1.0)  # Starting top position (with space for title)
        pic_width = Inches(2.5)  # Adjust image width
        pic_height = Inches(1.5)  # Adjust image height

        # Arrange images in rows and columns
        for i, image in enumerate(images):
            row = i  # Determine row
            left_offset = left + col * pic_width  - Inches(2)# X offset for the column
            top_offset = top + row * pic_height  # Y offset for the row
            slide.shapes.add_picture(image, left_offset, top_offset, pic_width, pic_height)

    for parameter, hg_lg_dict in images.items():
        col = 0
        # Create the slide layout with custom title positioning
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Empty slide layout
        
        # Title text
        slide_title = f"{parameter}"  

        # Add a textbox for the title and manually set the position
        title_shape = slide.shapes.add_textbox(Inches(5.5), Inches(0.2), Inches(9), Inches(0.75))
        text_frame = title_shape.text_frame
        text_frame.text = slide_title

        # Set font size for better readability
        for p in text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Inches(0.56)

        for hg_lg, suffix_dict in hg_lg_dict.items():
            for suffix, imgs in suffix_dict.items():
                if imgs:
                    col = col + 1 
                    # Add images to the next column
                    add_images_to_slide(slide, imgs, col)

    # Save the PowerPoint presentation
    prs.save(output_pptx)
    print(f"Presentation saved as {output_pptx}")  # Debugging: confirm saving

# Main function to drive the process
def main(csv_file, image_folder, output_pptx):
    parameters = load_parameters(csv_file)
    images = load_images(image_folder, parameters)
    create_presentation(images, output_pptx)

# Run the program
csv_file = 'parameter.csv'  # Excel file containing the parameter names
image_folder = '../BNL_Tray1_Tray4_Tray2_tray3_674/rstst/'  # Folder containing the images
output_pptx = 'output_presentation.pptx'  # Name of the output PowerPoint file

main(csv_file, image_folder, output_pptx)
